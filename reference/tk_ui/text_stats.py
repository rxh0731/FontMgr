# text_stats.py — 文字统计主窗口内嵌页面

import csv 
import json 
import os 
import re 
import shutil 
import tempfile 
import threading 
import tkinter as tk 
import unicodedata 
from tkinter import filedialog ,ttk 
from typing import Any ,Callable ,Optional 

import numpy as np 
from PIL import Image ,ImageGrab ,ImageSequence 
from pypinyin import lazy_pinyin 

import config 
from ui import theme 
from ui .widgets .custom_dialog import show_error ,show_info ,show_warning 
from ui .widgets .dark_helpers import apply_dark_titlebar ,set_window_icon 


class TextStatsPage (tk .Frame ):
    """经文正文提取、不重复汉字统计和字库缺字分析。"""

    _ocr_engine :Any =None 
    _image_extensions ={".png",".jpg",".jpeg",".bmp",".tif",".tiff",".webp"}

    def __init__ (self ,parent :tk .Widget ,on_close :Optional [Callable [[],None ]]=None )->None :
        super ().__init__ (parent ,bg =theme .BG_MAIN )
        self .pack (fill =tk .BOTH ,expand =True )
        self ._on_close =on_close 
        self ._file_path :Optional [str ]=None 
        self ._session =object ()
        self ._task_id =0 
        self ._refresh_job :Optional [str ]=None 
        self ._all_characters :list [str ]=[]
        self ._missing_characters :list [str ]=[]
        self ._progress_window :Optional [tk .Toplevel ]=None 
        self ._build_ui ()
        self .bind_all ("<Escape>",self ._close_page ,add ="+")

    def _close_page (self ,_event :Any =None )->str :
        """关闭文字统计并返回进入前页面。"""
        if self ._on_close :
            self ._on_close ()
        return "break"

    def destroy (self )->None :
        self ._session =object ()
        if self ._refresh_job :
            try :
                self .after_cancel (self ._refresh_job )
            except tk .TclError :
                pass 
        self ._close_progress_window ()
        try :
            self .unbind_all ("<Escape>")
        except tk .TclError :
            pass 
        super ().destroy ()

    def _build_ui (self )->None :
        top_bar =tk .Frame (self ,bg =theme .BG_PANEL )
        top_bar .pack (fill =tk .X )
        theme .make_label (top_bar ,"文字统计",font =theme .FONT_TITLE ).pack (side =tk .LEFT ,padx =10 ,pady =8 )
        if self ._on_close :
            theme .make_button (top_bar ,"返回首页",command =self ._close_page ).pack (side =tk .RIGHT ,padx =10 ,pady =8 )

        selection_frame =tk .Frame (self ,bg =theme .BG_MAIN )
        selection_frame .pack (fill =tk .X ,padx =15 ,pady =(12 ,4 ))
        selection_frame .grid_columnconfigure (1 ,weight =1 )
        selection_frame .grid_columnconfigure (3 ,minsize =80 )

        theme .make_label (selection_frame ,"经文文本：",width =10 ,anchor ="e").grid (row =0 ,column =0 ,sticky ="e")
        self ._path_var =tk .StringVar ()
        tk .Entry (
        selection_frame ,textvariable =self ._path_var ,state ="readonly",
        readonlybackground =theme .BG_INPUT ,fg =theme .FG_PRIMARY ,relief =tk .FLAT ,
        ).grid (row =0 ,column =1 ,sticky ="ew",padx =4 ,ipady =4 )
        self ._select_button =theme .make_button (selection_frame ,"选择文件",command =self ._select_files )
        self ._select_button .grid (row =0 ,column =2 ,padx =2 )
        theme .make_button (selection_frame ,"文字统计",accent =True ,command =self ._run_statistics ).grid (row =0 ,column =4 ,padx =2 )

        theme .make_label (selection_frame ,"字库目录：",width =10 ,anchor ="e").grid (row =1 ,column =0 ,sticky ="e",pady =(6 ,0 ))
        self ._font_dir_var =tk .StringVar ()
        tk .Entry (
        selection_frame ,textvariable =self ._font_dir_var ,state ="readonly",
        readonlybackground =theme .BG_INPUT ,fg =theme .FG_PRIMARY ,relief =tk .FLAT ,
        ).grid (row =1 ,column =1 ,sticky ="ew",padx =4 ,pady =(6 ,0 ),ipady =4 )
        theme .make_button (selection_frame ,"选择目录",command =self ._select_font_directory ).grid (row =1 ,column =2 ,padx =2 ,pady =(6 ,0 ))

        self ._info_var =tk .StringVar (value ="尚未选择文件")
        theme .make_label (self ,textvariable =self ._info_var ,fg =theme .FG_MUTED ,anchor ="w").pack (
        fill =tk .X ,padx =15 ,pady =(6 ,4 )
        )

        main_frame =tk .Frame (self ,bg =theme .BG_MAIN )
        main_frame .pack (fill =tk .BOTH ,expand =True ,padx =15 ,pady =(0 ,12 ))
        main_frame .grid_rowconfigure (0 ,weight =1 ,uniform ="文字统计")
        main_frame .grid_rowconfigure (1 ,weight =1 ,uniform ="文字统计")
        main_frame .grid_columnconfigure (0 ,weight =1 )

        content_frame =tk .Frame (main_frame ,bg =theme .BG_MAIN )
        content_frame .grid (row =0 ,column =0 ,sticky ="nsew",pady =(0 ,3 ))
        theme .make_label (content_frame ,"经文内容（可输入或粘贴文字、文件和图片）：",font =theme .FONT_BOLD ,fg =theme .FG_MUTED ).pack (
        fill =tk .X ,pady =(0 ,2 )
        )
        self ._content_text =tk .Text (
        content_frame ,font =theme .FONT_NORMAL ,fg =theme .FG_PRIMARY ,bg =theme .BG_INPUT ,
        insertbackground =theme .FG_PRIMARY ,wrap =tk .WORD ,relief =tk .FLAT ,
        highlightthickness =0 ,padx =12 ,pady =8 ,undo =True ,
        )
        content_scrollbar =tk .Scrollbar (content_frame ,command =self ._content_text .yview )
        self ._content_text .configure (yscrollcommand =content_scrollbar .set )
        content_scrollbar .pack (side =tk .RIGHT ,fill =tk .Y )
        self ._content_text .pack (fill =tk .BOTH ,expand =True )
        self ._content_text .bind ("<<Modified>>",self ._on_content_modified )
        for key_binding in ("<Control-v>","<Control-V>","<Shift-Insert>"):
            self ._content_text .bind (key_binding ,self ._paste )

        results_frame =tk .Frame (main_frame ,bg =theme .BG_MAIN )
        results_frame .grid (row =1 ,column =0 ,sticky ="nsew",pady =(3 ,0 ))
        results_frame .grid_columnconfigure (0 ,weight =1 ,uniform ="结果")
        results_frame .grid_columnconfigure (1 ,weight =1 ,uniform ="结果")
        results_frame .grid_rowconfigure (1 ,weight =1 )

        all_header =tk .Frame (results_frame ,bg =theme .BG_MAIN )
        all_header .grid (row =0 ,column =0 ,sticky ="ew",padx =(0 ,2 ))
        theme .make_label (all_header ,"统计结果（不重复汉字，按拼音排序）：",font =theme .FONT_BOLD ,fg =theme .FG_MUTED ).pack (side =tk .LEFT )
        self ._count_var =tk .StringVar ()
        theme .make_label (all_header ,textvariable =self ._count_var ,fg =theme .FG_MUTED ).pack (side =tk .LEFT ,padx =6 )
        theme .make_button (all_header ,"全部文字导出",accent =True ,command =self ._export_all ).pack (side =tk .RIGHT )

        missing_header =tk .Frame (results_frame ,bg =theme .BG_MAIN )
        missing_header .grid (row =0 ,column =1 ,sticky ="ew",padx =(2 ,0 ))
        theme .make_label (missing_header ,"缺失文字：",font =theme .FONT_BOLD ,fg =theme .FG_MUTED ).pack (side =tk .LEFT )
        self ._missing_count_var =tk .StringVar ()
        theme .make_label (missing_header ,textvariable =self ._missing_count_var ,fg =theme .FG_MUTED ).pack (side =tk .LEFT ,padx =6 )
        theme .make_button (missing_header ,"缺失文字导出",accent =True ,command =self ._export_missing ).pack (side =tk .RIGHT )

        self ._all_results_text =self ._create_readonly_text (results_frame )
        self ._all_results_text .grid (row =1 ,column =0 ,sticky ="nsew",padx =(0 ,2 ))
        self ._missing_results_text =self ._create_readonly_text (results_frame )
        self ._missing_results_text .grid (row =1 ,column =1 ,sticky ="nsew",padx =(2 ,0 ))

    @staticmethod 
    def _create_readonly_text (parent :tk .Widget )->tk .Text :
        text_widget =tk .Text (
        parent ,font =theme .FONT_NORMAL ,fg =theme .FG_PRIMARY ,bg =theme .BG_INPUT ,
        wrap =tk .CHAR ,relief =tk .FLAT ,highlightthickness =0 ,padx =15 ,pady =8 ,
        )
        text_widget .configure (state =tk .DISABLED )
        return text_widget 

    @staticmethod 
    def _set_readonly_content (text_widget :tk .Text ,content :str )->None :
        text_widget .configure (state =tk .NORMAL )
        text_widget .delete ("1.0",tk .END )
        text_widget .insert ("1.0",content )
        text_widget .configure (state =tk .DISABLED )

    def _select_files (self )->None :
        pattern ="*.txt;*.rtf;*.doc;*.docx;*.pdf;*.md;*.html;*.htm;*.xml;*.csv;*.json;*.wps;*.wpt;*.et;*.ett;*.dps;*.dpt;*.xlsx;*.pptx;*.png;*.jpg;*.jpeg;*.bmp;*.tif;*.tiff;*.webp"
        paths =filedialog .askopenfilenames (
        parent =self .winfo_toplevel (),title ="选择经文文件或图片（可多选）",
        filetypes =[("支持的文件",pattern ),("图片文件","*.png;*.jpg;*.jpeg;*.bmp;*.tif;*.tiff;*.webp"),("所有文件","*.*")],
        )
        if paths :
            self ._start_file_task (list (paths ),replace =True )

    def _show_progress_window (self ,title :str ="正在处理文件")->None :
        self ._close_progress_window ()
        main_window =self .winfo_toplevel ()
        window =tk .Toplevel (main_window )
        self ._progress_window =window 
        window .withdraw ()
        window .title (title )
        window .configure (bg =theme .BG_PANEL )
        window .resizable (False ,False )
        window .transient (main_window )
        set_window_icon (window ,config .ICON_FILE )
        apply_dark_titlebar (window )
        window .protocol ("WM_DELETE_WINDOW",lambda :None )
        theme .make_label (window ,title ,font =theme .FONT_TITLE ).pack (pady =(18 ,10 ))
        self ._progress_label =theme .make_label (window ,"正在准备……",width =48 ,anchor ="w")
        self ._progress_label .pack (padx =24 )
        self ._progress_var =tk .DoubleVar (value =0 )
        ttk .Progressbar (window ,variable =self ._progress_var ,maximum =100 ,length =420 ).pack (padx =24 ,pady =(10 ,6 ))
        self ._progress_percent_label =theme .make_label (window ,"0%",fg =theme .FG_MUTED )
        self ._progress_percent_label .pack (pady =(0 ,16 ))
        window .update_idletasks ()
        width ,height =window .winfo_reqwidth (),window .winfo_reqheight ()
        x =max (0 ,(window .winfo_screenwidth ()-width )//2 )
        y =max (0 ,(window .winfo_screenheight ()-height )//2 )
        window .geometry (f"{width }x{height }+{x }+{y }")
        window .deiconify ()
        window .lift ()
        window .grab_set ()

    def _update_progress (self ,session :object ,task_id :int ,value :float ,description :str )->None :
        def apply_update ()->None :
            if session is not self ._session or task_id !=self ._task_id :
                return 
            if self ._progress_window is None or not self ._progress_window .winfo_exists ():
                return 
            bounded_value =max (0.0 ,min (100.0 ,float (value )))
            self ._progress_var .set (bounded_value )
            self ._progress_label .configure (text =description )
            self ._progress_percent_label .configure (text =f"{bounded_value :.0f}%")

        try :
            self .after (0 ,apply_update )
        except tk .TclError :
            pass 

    def _close_progress_window (self )->None :
        if self ._progress_window is not None :
            try :
                self ._progress_window .grab_release ()
                self ._progress_window .destroy ()
            except tk .TclError :
                pass 
        self ._progress_window =None 

    def _start_file_task (self ,paths :list [str ],replace :bool =True ,clipboard_image :Optional [Image .Image ]=None )->None :
        session =self ._session 
        self ._task_id +=1 
        task_id =self ._task_id 
        item_count =1 if clipboard_image is not None else len (paths )
        self ._info_var .set (f"正在后台处理 {item_count } 个文件……")
        self ._select_button .configure (state =tk .DISABLED ,text ="正在处理")
        self ._show_progress_window ("正在提取经文文字")

        def worker ()->None :
            values :list [str ]=[]
            report_list :list [str ]=[]
            failure_list :list [str ]=[]
            items =[("剪贴板图片",clipboard_image )]if clipboard_image is not None else [(path ,None )for path in paths ]
            total =max (1 ,len (items ))
            for index ,(path ,image )in enumerate (items ):
                name ="剪贴板图片"if image is not None else os .path .basename (path )
                base_progress ,progress_span =index /total *100 ,100 /total 
                self ._update_progress (session ,task_id ,base_progress ,f"正在处理：{name }")
                try :
                    def report_progress (ratio :float ,detail :str )->None :
                        self ._update_progress (session ,task_id ,base_progress +progress_span *max (0.0 ,min (1.0 ,ratio )),f"{name }：{detail }")
                    if image is not None :
                        text =self ._recognize_image (image ,report_progress )
                        report ="已识别剪贴板图片文字"
                    else :
                        text ,report =self ._read_file_text (path ,report_progress )
                    if text .strip ():
                        values .append (text )
                        report_list .append (f"{name }：{report }")
                    else :
                        failure_list .append (f"{name }：未识别或提取到正文")
                except Exception as error :
                    failure_list .append (f"{name }：{error }")
                self ._update_progress (session ,task_id ,(index +1 )/total *100 ,f"已处理：{name }")
            try :
                self .after (0 ,lambda :finish (values ,report_list ,failure_list ))
            except tk .TclError :
                pass 

        def finish (values :list [str ],report_list :list [str ],failure_list :list [str ])->None :
            if session is not self ._session or task_id !=self ._task_id or not self .winfo_exists ():
                return 
            self ._close_progress_window ()
            self ._select_button .configure (state =tk .NORMAL ,text ="选择文件")
            if values :
                combined_content ="\n\n".join (values )
                if replace :
                    self ._content_text .delete ("1.0",tk .END )
                    self ._content_text .insert ("1.0",combined_content )
                else :
                    if self ._content_text .tag_ranges (tk .SEL ):
                        self ._content_text .delete (tk .SEL_FIRST ,tk .SEL_LAST )
                    self ._content_text .insert (tk .INSERT ,combined_content )
                self ._file_path =paths [0 ]if len (paths )==1 and clipboard_image is None else None 
                self ._path_var .set (paths [0 ]if self ._file_path else f"已处理 {len (values )} 个文件或图片")
                self ._content_text .edit_modified (False )
                self ._refresh_info ()
                self ._run_statistics (silent =True )
            summary =[f"成功处理：{len (values )} 个文件或图片",*report_list ]
            if failure_list :
                summary .extend ([f"未能处理：{len (failure_list )} 个文件或图片",*failure_list ])
            show_info (self .winfo_toplevel (),"文件处理结果","\n".join (summary ))

        threading .Thread (target =worker ,daemon =True ).start ()

    def _paste (self ,_event :Any =None )->str :
        try :
            text =self .clipboard_get ()
        except tk .TclError :
            text =None 
        if text is not None :
            text =self ._clean_plain_text (text )
            if self ._content_text .tag_ranges (tk .SEL ):
                self ._content_text .delete (tk .SEL_FIRST ,tk .SEL_LAST )
            self ._content_text .insert (tk .INSERT ,text )
            return "break"
        try :
            clipboard_content =ImageGrab .grabclipboard ()
        except Exception :
            clipboard_content =None 
        if isinstance (clipboard_content ,list ):
            supported_paths =[path for path in clipboard_content if os .path .isfile (path )]
            if supported_paths :
                self ._start_file_task (supported_paths ,replace =False )
                return "break"
        if isinstance (clipboard_content ,Image .Image ):
            self ._start_file_task ([],replace =False ,clipboard_image =clipboard_content .copy ())
            return "break"
        show_warning (self .winfo_toplevel (),"粘贴结果","剪贴板中没有可用于文字统计的内容。")
        return "break"

    def _on_content_modified (self ,_event :Any =None )->None :
        if not self ._content_text .edit_modified ():
            return 
        self ._content_text .edit_modified (False )
        if self ._file_path is not None :
            self ._file_path =None 
            self ._path_var .set ("")
        if self ._refresh_job :
            try :
                self .after_cancel (self ._refresh_job )
            except tk .TclError :
                pass 
        session =self ._session 
        self ._refresh_job =self .after (400 ,lambda :self ._refresh_after_edit (session ))

    def _refresh_after_edit (self ,session :object )->None :
        self ._refresh_job =None 
        if session is not self ._session :
            return 
        self ._refresh_info ()
        if self ._content_text .get ("1.0","end-1c").strip ():
            self ._run_statistics (silent =True )

    def _refresh_info (self )->None :
        content =self ._content_text .get ("1.0","end-1c")
        if not content .strip ():
            self ._info_var .set ("尚未载入文本")
            return 
        stats =self ._analyze_text (content )
        prefix =f"文件：{os .path .basename (self ._file_path )}  |  "if self ._file_path else ""
        self ._info_var .set (
        f"{prefix }总字符：{stats ['总字数']}  |  汉字：{stats ['中文文字数']}  |  "
        f"英文词：{stats ['英文文字数']}  |  标点：{stats ['标点数']}  |  空白：{stats ['空格数']}"
        )

    def _select_font_directory (self )->None :
        """打开系统目录选择器，并确保它显示在主窗口前方。"""
        main_window =self .winfo_toplevel ()
        main_window .update_idletasks ()
        directory =""
        try :
            main_window .attributes ("-topmost",True )
            directory =filedialog .askdirectory (
            parent =main_window ,
            title ="选择字库目录（含文字图片的文件夹）",
            mustexist =True ,
            )
        finally :
            try :
                main_window .attributes ("-topmost",False )
                main_window .lift ()
                main_window .focus_force ()
            except tk .TclError :
                pass 
        if directory :
            self ._font_dir_var .set (directory )
            self ._analyze_missing_characters ()

    @staticmethod 
    def _is_chinese_character (character :str )->bool :
        if not character :
            return False 
        codepoint =ord (character )
        return (
        0x3400 <=codepoint <=0x4DBF or 0x4E00 <=codepoint <=0x9FFF or 
        0xF900 <=codepoint <=0xFAFF or 0x20000 <=codepoint <=0x2EE5F or 
        0x2F800 <=codepoint <=0x2FA1F or 0x30000 <=codepoint <=0x323AF 
        )

    @staticmethod 
    def _pinyin_sort_key (character :str )->tuple [str ,str ]:
        try :
            pronunciation =lazy_pinyin (character )
            return (pronunciation [0 ]if pronunciation else character ,character )
        except Exception :
            return (character ,character )

    def _run_statistics (self ,silent :bool =False )->None :
        content =self ._content_text .get ("1.0","end-1c")
        if not content .strip ():
            if not silent :
                show_warning (self .winfo_toplevel (),"提示","请先载入文件或在内容窗口粘贴文本。")
            return 
        self ._all_characters =sorted ({character for character in content if self ._is_chinese_character (character )},key =self ._pinyin_sort_key )
        self ._count_var .set (f"共 {len (self ._all_characters )} 个不重复汉字")
        self ._set_readonly_content (self ._all_results_text ," ".join (self ._all_characters ))
        self ._analyze_missing_characters ()

    def _analyze_missing_characters (self )->None :
        if not self ._all_characters :
            return 
        font_directory =self ._font_dir_var .get ().strip ()
        if not font_directory or not os .path .isdir (font_directory ):
            self ._missing_count_var .set ("请选择字库目录")
            self ._missing_characters =[]
            missing_text ="（选择字库目录后分析缺失文字）"
        else :
            image_extensions =self ._image_extensions |{".gif",".psd",".tga",".ico"}
            existing_characters :set [str ]=set ()
            invalid_name_count =0 
            try :
                for _ ,_ ,file_names in os .walk (font_directory ):
                    for file_name in file_names :
                        if os .path .splitext (file_name )[1 ].lower ()not in image_extensions :
                            continue 
                        stem =os .path .splitext (file_name )[0 ]
                        candidate_character =stem .rsplit ("-",1 )[0 ]if "-"in stem else stem 
                        if len (candidate_character )==1 and self ._is_chinese_character (candidate_character ):
                            existing_characters .add (candidate_character )
                        else :
                            invalid_name_count +=1 
                self ._missing_characters =sorted (set (self ._all_characters )-existing_characters ,key =self ._pinyin_sort_key )
                suffix =f"，忽略 {invalid_name_count } 个名称异常文件"if invalid_name_count else ""
                self ._missing_count_var .set (f"共 {len (self ._missing_characters )} 个缺失{suffix }")
                missing_text =" ".join (self ._missing_characters )if self ._missing_characters else "（无缺失文字）"
            except Exception as error :
                self ._missing_characters =[]
                self ._missing_count_var .set ("字库目录读取失败")
                missing_text =f"无法读取字库目录：{error }"
        self ._set_readonly_content (self ._missing_results_text ,missing_text )

    def _export_all (self )->None :
        self ._export_characters (self ._all_characters ,"导出全部文字","统计结果.txt")

    def _export_missing (self )->None :
        self ._export_characters (self ._missing_characters ,"导出缺失文字","缺失文字.txt")

    def _export_characters (self ,characters :list [str ],title :str ,default_name :str )->None :
        if not characters :
            show_warning (self .winfo_toplevel (),"提示","当前没有可导出的文字。")
            return 
        path =filedialog .asksaveasfilename (
        parent =self .winfo_toplevel (),title =title ,defaultextension =".txt",
        initialfile =default_name ,filetypes =[("文本文件","*.txt")],
        )
        if not path :
            return 
        try :
            with open (path ,"w",encoding ="utf-8")as file :
                file .write (" ".join (characters ))
            show_info (self .winfo_toplevel (),"导出成功",f"已导出到：\n{path }")
        except Exception as error :
            show_error (self .winfo_toplevel (),"导出失败",f"无法导出：{error }")

    @staticmethod 
    def _clean_plain_text (text :str )->str :
        text =str (text ).replace ("\ufeff","").replace ("\u200b","").replace ("\x00","")
        text =text .replace ("\r\n","\n").replace ("\r","\n")
        text ="".join (character for character in text if character in "\n\t"or ord (character )>=32 )
        return re .sub (r"\n{4,}","\n\n\n",text ).strip ()

    @staticmethod 
    def _read_encoded_text (path :str )->str :
        with open (path ,"rb")as file :
            raw_content =file .read ()
        for codepoint in ("utf-8-sig","utf-16","utf-32","gb18030"):
            try :
                return raw_content .decode (codepoint )
            except UnicodeError :
                pass 
        try :
            from charset_normalizer import from_bytes 
            best_result =from_bytes (raw_content ).best ()
            if best_result is not None :
                return str (best_result )
        except ImportError :
            pass 
        raise ValueError ("无法可靠识别文本编码")

    @staticmethod 
    def _read_word_document (path :str )->str :
        try :
            from docx import Document 
            from docx .oxml .table import CT_Tbl 
            from docx .oxml .text .paragraph import CT_P 
            from docx .table import Table 
            from docx .text .paragraph import Paragraph 
        except ImportError as error :
            raise ValueError ("缺少文字文档读取组件")from error 
        document =Document (path )
        lines :list [str ]=[]
        for child in document .element .body .iterchildren ():
            if isinstance (child ,CT_P ):
                content =Paragraph (child ,document ).text .strip ()
                if content :
                    lines .append (content )
            elif isinstance (child ,CT_Tbl ):
                for row in Table (child ,document ).rows :
                    values =[cell .text .strip ().replace ("\n"," ")for cell in row .cells ]
                    if any (values ):
                        lines .append ("\t".join (values ))
        for section in document .sections :
            for region in (section .header ,section .footer ):
                lines .extend (paragraph .text .strip ()for paragraph in region .paragraphs if paragraph .text .strip ())
        return "\n".join (lines )

    @staticmethod 
    def _convert_office_document (path :str ,target_extension :str )->str :
        try :
            import pythoncom 
            import win32com .client 
        except ImportError as error :
            raise ValueError ("缺少本机办公文档转换组件")from error 
        if target_extension ==".docx":
            application_names ,document_type ,save_format =("Kwps.Application","Word.Application"),"文字",16 
        elif target_extension ==".xlsx":
            application_names ,document_type ,save_format =("Ket.Application","Excel.Application"),"表格",51 
        else :
            application_names ,document_type ,save_format =("Kwpp.Application","PowerPoint.Application"),"演示",24 
        pythoncom .CoInitialize ()
        try :
            for application_name in application_names :
                application =document =None 
                temp_directory =tempfile .mkdtemp (prefix ="字库编辑_文字统计_")
                output_path =os .path .join (temp_directory ,os .path .splitext (os .path .basename (path ))[0 ]+target_extension )
                try :
                    application =win32com .client .DispatchEx (application_name )
                    application .Visible =False 
                    if target_extension ==".docx":
                        document =application .Documents .Open (os .path .abspath (path ),ReadOnly =True )
                        document .SaveAs2 (output_path ,FileFormat =save_format )
                    elif target_extension ==".xlsx":
                        document =application .Workbooks .Open (os .path .abspath (path ),ReadOnly =True )
                        document .SaveAs (output_path ,FileFormat =save_format )
                    else :
                        document =application .Presentations .Open (os .path .abspath (path ),WithWindow =False )
                        document .SaveAs (output_path ,save_format )
                    if os .path .isfile (output_path ):
                        return output_path 
                except Exception :
                    shutil .rmtree (temp_directory ,ignore_errors =True )
                finally :
                    try :
                        if document is not None :
                            document .Close (False )
                    except Exception :
                        pass 
                    try :
                        if application is not None :
                            application .Quit ()
                    except Exception :
                        pass 
        finally :
            pythoncom .CoUninitialize ()
        raise ValueError (f"无法调用本机 WPS 或 Office 转换{document_type }文档")

    @staticmethod 
    def _arrange_ocr_order (texts :list [str ],boxes :Any )->tuple [str ,str ]:
        if not texts or boxes is None or len (texts )!=len (boxes ):
            return "\n".join (texts or []),"横排"
        entries :list [dict [str ,Any ]]=[]
        for text ,box in zip (texts ,boxes ):
            coordinates =np .asarray (box ,dtype =float )
            if coordinates .shape !=(4 ,2 )or not str (text ).strip ():
                continue 
            left ,right =float (coordinates [:,0 ].min ()),float (coordinates [:,0 ].max ())
            top ,bottom =float (coordinates [:,1 ].min ()),float (coordinates [:,1 ].max ())
            entries .append ({"文字":str (text ).strip (),"横心":(left +right )/2 ,"纵心":(top +bottom )/2 ,
            "宽":max (1.0 ,right -left ),"高":max (1.0 ,bottom -top )})
        if not entries :
            return "","横排"
        median_width =float (np .median ([entry ["宽"]for entry in entries ]))
        median_height =float (np .median ([entry ["高"]for entry in entries ]))
        horizontal_evidence =sum (max (0.0 ,entry ["宽"]/entry ["高"]-1.25 )for entry in entries )
        vertical_evidence =sum (max (0.0 ,entry ["高"]/entry ["宽"]-1.25 )for entry in entries )
        horizontal_alignment =vertical_alignment =0 
        for task_id ,current in enumerate (entries ):
            for other in entries [task_id +1 :]:
                if abs (current ["纵心"]-other ["纵心"])<=median_height *0.65 :
                    horizontal_alignment +=1 
                if abs (current ["横心"]-other ["横心"])<=median_width *0.65 :
                    vertical_alignment +=1 
        is_vertical =vertical_evidence *2 +vertical_alignment >(horizontal_evidence *2 +horizontal_alignment )*1.2 and (vertical_evidence >0 or vertical_alignment >=2 )
        groups :list [list [dict [str ,Any ]]]=[]
        sort_key ="横心"if is_vertical else "纵心"
        tolerance =(median_width if is_vertical else median_height )*0.7 
        for current in sorted (entries ,key =lambda entry :entry [sort_key ],reverse =is_vertical ):
            best_group =None 
            minimum_distance =float ("inf")
            for group in groups :
                group_center =sum (entry [sort_key ]for entry in group )/len (group )
                distance =abs (current [sort_key ]-group_center )
                if distance <=tolerance and distance <minimum_distance :
                    best_group ,minimum_distance =group ,distance 
            if best_group is None :
                groups .append ([current ])
            else :
                best_group .append (current )
        if is_vertical :
            groups .sort (key =lambda group :sum (entry ["横心"]for entry in group )/len (group ),reverse =True )
            ordered_lines =["".join (entry ["文字"]for entry in sorted (group ,key =lambda entry :entry ["纵心"]))for group in groups ]
            return "\n".join (ordered_lines ),"竖排（从右向左、从上到下）"
        groups .sort (key =lambda group :sum (entry ["纵心"]for entry in group )/len (group ))
        ordered_lines =["".join (entry ["文字"]for entry in sorted (group ,key =lambda entry :entry ["横心"]))for group in groups ]
        return "\n".join (ordered_lines ),"横排（从上到下、从左向右）"

    @classmethod 
    def _recognize_image (cls ,image :Image .Image ,progress :Optional [Callable [[float ,str ],None ]]=None )->str :
        if progress :
            progress (0.05 ,"正在加载图片识别引擎")
        try :
            from rapidocr import RapidOCR 
            if cls ._ocr_engine is None :
                cls ._ocr_engine =RapidOCR ()
            if progress :
                progress (0.25 ,"正在分析图片文字")
            result =cls ._ocr_engine (np .array (image .convert ("RGB")))
            texts =list (getattr (result ,"txts",None )or [])
            boxes =getattr (result ,"boxes",None )
            text ,layout =cls ._arrange_ocr_order (texts ,boxes )
            if progress :
                progress (1.0 ,f"识别完成，已判断为{layout }")
            return cls ._clean_plain_text (text )
        except ImportError :
            try :
                import pytesseract 
                if progress :
                    progress (0.25 ,"正在使用后备引擎识别")
                text =pytesseract .image_to_string (image ,lang ="chi_sim+chi_tra")
                if progress :
                    progress (1.0 ,"图片识别完成")
                return cls ._clean_plain_text (text )
            except Exception as error :
                raise ValueError (f"图片识别引擎不可用：{error }")from error 

    @classmethod 
    def _recognize_image_file (cls ,path :str ,progress :Optional [Callable [[float ,str ],None ]]=None )->str :
        with Image .open (path )as image :
            pages =[page .copy ()for page in ImageSequence .Iterator (image )]
        total =max (1 ,len (pages ))
        text_list :list [str ]=[]
        for index ,page in enumerate (pages ):
            def page_progress (ratio :float ,detail :str )->None :
                if progress :
                    progress ((index +ratio )/total ,f"第 {index +1 }/{total } 页：{detail }")
            text =cls ._recognize_image (page ,page_progress )
            if text :
                text_list .append (text )
        return cls ._clean_plain_text ("\n\n".join (text_list ))

    @classmethod 
    def _read_file_text (cls ,path :str ,progress :Optional [Callable [[float ,str ],None ]]=None )->tuple [str ,str ]:
        extension =os .path .splitext (path )[1 ].lower ()
        if extension in cls ._image_extensions :
            return cls ._recognize_image_file (path ,progress ),"已使用本地识别引擎提取图片文字"
        if extension in {".txt",".md"}:
            return cls ._clean_plain_text (cls ._read_encoded_text (path )),"已按纯文本提取"
        if extension ==".rtf":
            from striprtf .striprtf import rtf_to_text 
            return cls ._clean_plain_text (rtf_to_text (cls ._read_encoded_text (path ))),"已去除富文本格式"
        if extension in {".html",".htm"}:
            from bs4 import BeautifulSoup 
            document =BeautifulSoup (cls ._read_encoded_text (path ),"html.parser")
            for node in document (["script","style","noscript","template"]):
                node .decompose ()
            return cls ._clean_plain_text (document .get_text ("\n")),"已提取网页可见文字"
        if extension ==".xml":
            source_text =cls ._read_encoded_text (path )
            if "<!DOCTYPE"in source_text .upper ():
                raise ValueError ("不读取包含外部实体声明的 XML")
            import xml .etree .ElementTree as ET 
            root =ET .fromstring (source_text )
            text ="\n".join (content .strip ()for content in root .itertext ()if content .strip ())
            return cls ._clean_plain_text (text ),"已提取 XML 节点正文"
        if extension ==".csv":
            source_text =cls ._read_encoded_text (path )
            try :
                dialect =csv .Sniffer ().sniff (source_text [:8192 ],delimiters =",\t;|")
            except csv .Error :
                dialect =csv .excel 
            text ="\n".join ("\t".join (content .strip ()for content in row )for row in csv .reader (source_text .splitlines (),dialect ))
            return cls ._clean_plain_text (text ),"已提取表格单元格"
        if extension ==".json":
            strings :list [str ]=[]
            def collect (content :Any )->None :
                if isinstance (content ,str ):
                    strings .append (content )
                elif isinstance (content ,list ):
                    for child in content :
                        collect (child )
                elif isinstance (content ,dict ):
                    for child in content .values ():
                        collect (child )
            collect (json .loads (cls ._read_encoded_text (path )))
            return cls ._clean_plain_text ("\n".join (strings )),"已提取字符串值，未统计字段名"
        if extension ==".docx":
            return cls ._clean_plain_text (cls ._read_word_document (path )),"已提取段落、表格、页眉和页脚"
        if extension ==".xlsx":
            from openpyxl import load_workbook 
            workbook =load_workbook (path ,read_only =True ,data_only =True )
            lines :list [str ]=[]
            try :
                for worksheet in workbook .worksheets :
                    for row in worksheet .iter_rows (values_only =True ):
                        values =[""if content is None else str (content )for content in row ]
                        if any (values ):
                            lines .append ("\t".join (values ))
            finally :
                workbook .close ()
            return cls ._clean_plain_text ("\n".join (lines )),"已提取所有工作表单元格"
        if extension ==".pptx":
            from pptx import Presentation 
            lines =[]
            for slide in Presentation (path ).slides :
                for shape in slide .shapes :
                    if getattr (shape ,"has_text_frame",False )and shape .text .strip ():
                        lines .append (shape .text .strip ())
                    if getattr (shape ,"has_table",False ):
                        for row in shape .table .rows :
                            values =[cell .text .strip ()for cell in row .cells ]
                            if any (values ):
                                lines .append ("\t".join (values ))
            return cls ._clean_plain_text ("\n".join (lines )),"已提取幻灯片和表格文字"
        if extension ==".pdf":
            import fitz 
            document =fitz .open (path )
            text_list :list [str ]=[]
            try :
                page_count =len (document )
                for index ,page in enumerate (document ):
                    if progress :
                        progress (index /max (1 ,page_count ),f"正在读取第 {index +1 }/{page_count } 页")
                    page_text =page .get_text ("text",sort =True )
                    if sum (cls ._is_chinese_character (character )for character in page_text )<5 :
                        pixmap =page .get_pixmap (matrix =fitz .Matrix (2 ,2 ),alpha =False )
                        image =Image .frombytes ("RGB",(pixmap .width ,pixmap .height ),pixmap .samples )
                        def page_progress (ratio :float ,detail :str )->None :
                            if progress :
                                progress ((index +ratio )/max (1 ,page_count ),f"第 {index +1 }/{page_count } 页：{detail }")
                        page_text =cls ._recognize_image (image ,page_progress )
                    if page_text .strip ():
                        text_list .append (page_text )
                if progress :
                    progress (1.0 ,"PDF 处理完成")
            finally :
                document .close ()
            text =cls ._clean_plain_text ("\n\n".join (text_list ))
            if not text :
                raise ValueError ("PDF 文本层和扫描页面均未识别出文字")
            return text ,"已提取 PDF 文本层，扫描页面已自动识别"
        if extension in {".doc",".wps",".wpt"}:
            converted_path =cls ._convert_office_document (path ,".docx")
            text ,_ =cls ._read_file_text (converted_path )
            return text ,"已由 WPS 或 Office 后台转换并提取文字文档"
        if extension in {".et",".ett"}:
            converted_path =cls ._convert_office_document (path ,".xlsx")
            text ,_ =cls ._read_file_text (converted_path )
            return text ,"已由 WPS 后台转换并提取表格单元格"
        if extension in {".dps",".dpt"}:
            converted_path =cls ._convert_office_document (path ,".pptx")
            text ,_ =cls ._read_file_text (converted_path )
            return text ,"已由 WPS 后台转换并提取演示文稿文字"
        raise ValueError (f"暂不支持该文件格式：{extension or '无扩展名'}")

    @classmethod 
    def _analyze_text (cls ,text :str )->dict [str ,int ]:
        return {
        "总字数":len (text ),
        "标点数":sum (unicodedata .category (character ).startswith ("P")for character in text ),
        "空格数":sum (character .isspace ()for character in text ),
        "中文文字数":sum (cls ._is_chinese_character (character )for character in text ),
        "英文文字数":len (re .findall (r"[a-zA-Z]+",text )),
        }
