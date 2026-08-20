"""经文正文文件的只读文字提取服务。"""

from __future__ import annotations

import csv
import json
import os
import re
import xml.etree.ElementTree as ET
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any


MAX_SOURCE_FILE_BYTES = 512 * 1024 * 1024
MAX_TEXT_SOURCE_BYTES = 64 * 1024 * 1024
MAX_EXTRACTED_CHARACTERS = 10_000_000
MAX_ARCHIVE_MEMBERS = 20_000
MAX_ARCHIVE_ENTRY_BYTES = 128 * 1024 * 1024
MAX_ARCHIVE_TOTAL_BYTES = 512 * 1024 * 1024
MAX_ARCHIVE_COMPRESSION_RATIO = 1000
MAX_PDF_PAGES = 10_000
MAX_SPREADSHEET_CELLS = 5_000_000
MAX_PRESENTATION_SLIDES = 10_000


PLAIN_TEXT_EXTENSIONS = frozenset(
    {
        ".txt",
        ".md",
        ".markdown",
        ".log",
        ".ini",
        ".yaml",
        ".yml",
    }
)
SUPPORTED_DOCUMENT_EXTENSIONS = frozenset(
    {
        *PLAIN_TEXT_EXTENSIONS,
        ".rtf",
        ".html",
        ".htm",
        ".xml",
        ".csv",
        ".tsv",
        ".json",
        ".doc",
        ".docx",
        ".wps",
        ".wpt",
        ".pdf",
        ".xls",
        ".xlsx",
        ".et",
        ".ett",
        ".ppt",
        ".pptx",
        ".dps",
        ".dpt",
        ".odt",
        ".ods",
        ".odp",
    }
)


@dataclass(frozen=True, slots=True)
class ScriptureTextResult:
    text: str
    source_name: str
    detail: str


def scripture_document_filter() -> str:
    extensions = " ".join(f"*{suffix}" for suffix in sorted(SUPPORTED_DOCUMENT_EXTENSIONS))
    return (
        f"支持的文档（{extensions}）;;"
        "文字文档（*.txt *.md *.markdown *.rtf *.doc *.docx *.wps *.wpt *.odt）;;"
        "PDF 文档（*.pdf）;;"
        "网页与结构化文本（*.html *.htm *.xml *.json *.yaml *.yml *.csv *.tsv）;;"
        "表格与演示文稿（*.xls *.xlsx *.et *.ett *.ppt *.pptx *.dps *.dpt *.ods *.odp）;;"
        "所有文件（*.*）"
    )


def load_scripture_text(path: str) -> ScriptureTextResult:
    """只读提取文档文字，不修改源文件。"""

    normalized = os.path.abspath(os.fspath(path))
    if not os.path.isfile(normalized) or os.path.islink(normalized):
        raise ValueError("所选文件不存在，或不是可安全读取的普通文件。")
    extension = Path(normalized).suffix.lower()
    if extension not in SUPPORTED_DOCUMENT_EXTENSIONS:
        raise ValueError(f"暂不支持该文件格式：{extension or '无扩展名'}")
    source_size = os.path.getsize(normalized)
    if source_size > MAX_SOURCE_FILE_BYTES:
        raise ValueError(
            f"文档大小超过安全上限（{MAX_SOURCE_FILE_BYTES // 1024 // 1024} MiB）。"
        )
    if extension in {
        *PLAIN_TEXT_EXTENSIONS,
        ".rtf",
        ".html",
        ".htm",
        ".xml",
        ".csv",
        ".tsv",
        ".json",
    } and source_size > MAX_TEXT_SOURCE_BYTES:
        raise ValueError(
            f"文本类文档大小超过安全上限（{MAX_TEXT_SOURCE_BYTES // 1024 // 1024} MiB）。"
        )
    if extension in {".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"}:
        _validate_zip_archive(normalized)

    if extension in PLAIN_TEXT_EXTENSIONS:
        text, detail = _read_encoded_text(normalized), "已按纯文本读取"
    elif extension == ".rtf":
        from striprtf.striprtf import rtf_to_text

        text, detail = rtf_to_text(_read_encoded_text(normalized)), "已去除富文本格式"
    elif extension in {".html", ".htm"}:
        from bs4 import BeautifulSoup

        document = BeautifulSoup(_read_encoded_text(normalized), "html.parser")
        for node in document(["script", "style", "noscript", "template"]):
            node.decompose()
        text, detail = document.get_text("\n"), "已提取网页可见文字"
    elif extension == ".xml":
        text, detail = _read_xml(normalized), "已提取 XML 节点正文"
    elif extension in {".csv", ".tsv"}:
        text, detail = _read_delimited(normalized, extension), "已提取表格单元格"
    elif extension == ".json":
        text, detail = _read_json(normalized), "已提取字符串值"
    elif extension == ".docx":
        text, detail = _read_docx(normalized), "已提取段落、表格、页眉和页脚"
    elif extension == ".pdf":
        text, detail = _read_pdf(normalized), "已提取 PDF 文本层"
    elif extension == ".xlsx":
        text, detail = _read_xlsx(normalized), "已提取所有工作表单元格"
    elif extension == ".pptx":
        text, detail = _read_pptx(normalized), "已提取幻灯片和表格文字"
    elif extension in {".odt", ".ods", ".odp"}:
        text, detail = _read_open_document(normalized), "已提取开放文档正文"
    else:
        text, detail = _read_legacy_office(normalized, extension)

    if len(text) > MAX_EXTRACTED_CHARACTERS:
        raise ValueError(
            f"文档提取文字超过安全上限（{MAX_EXTRACTED_CHARACTERS:,} 个字符）。"
        )
    cleaned = _clean_plain_text(text)
    if not cleaned:
        raise ValueError("文档中没有提取到可用文字。")
    return ScriptureTextResult(cleaned, os.path.basename(normalized), detail)


def _clean_plain_text(text: object) -> str:
    value = str(text or "").replace("\ufeff", "").replace("\u200b", "")
    value = value.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    value = "".join(character for character in value if character in "\n\t" or ord(character) >= 32)
    return re.sub(r"\n{4,}", "\n\n\n", value).strip("\n")


def _read_encoded_text(path: str) -> str:
    with open(path, "rb") as stream:
        raw = stream.read()
    for encoding in ("utf-8-sig", "utf-16", "utf-32", "gb18030"):
        try:
            return raw.decode(encoding)
        except UnicodeError:
            continue
    from charset_normalizer import from_bytes

    result = from_bytes(raw).best()
    if result is None:
        raise ValueError("无法可靠识别文本编码。")
    return str(result)


def _validate_zip_archive(path: str) -> None:
    """限制 Office/OpenDocument 压缩包的展开规模和压缩比。"""

    try:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
    except (OSError, zipfile.BadZipFile) as exc:
        raise ValueError("文档压缩结构损坏，无法安全读取。") from exc
    if len(members) > MAX_ARCHIVE_MEMBERS:
        raise ValueError(f"文档内部文件数量超过安全上限（{MAX_ARCHIVE_MEMBERS:,} 项）。")
    total_size = 0
    for member in members:
        if member.file_size > MAX_ARCHIVE_ENTRY_BYTES:
            raise ValueError(
                f"文档内部文件“{member.filename}”展开后超过安全上限。"
            )
        total_size += member.file_size
        if total_size > MAX_ARCHIVE_TOTAL_BYTES:
            raise ValueError("文档解压后的总大小超过安全上限。")
        if (
            member.file_size > 0
            and member.compress_size == 0
        ) or (
            member.compress_size > 0
            and member.file_size / member.compress_size > MAX_ARCHIVE_COMPRESSION_RATIO
        ):
            raise ValueError(
                f"文档内部文件“{member.filename}”压缩比异常，已停止读取。"
            )


def _read_xml(path: str) -> str:
    source = _read_encoded_text(path)
    if "<!DOCTYPE" in source.upper():
        raise ValueError("为安全起见，不读取包含外部实体声明的 XML。")
    root = ET.fromstring(source)
    return "\n".join(value.strip() for value in root.itertext() if value.strip())


def _read_delimited(path: str, extension: str) -> str:
    source = _read_encoded_text(path)
    if extension == ".tsv":
        dialect: Any = csv.excel_tab
    else:
        try:
            dialect = csv.Sniffer().sniff(source[:8192], delimiters=",\t;|")
        except csv.Error:
            dialect = csv.excel
    return "\n".join(
        "\t".join(value.strip() for value in row)
        for row in csv.reader(source.splitlines(), dialect)
    )


def _read_json(path: str) -> str:
    values: list[str] = []

    def collect(value: Any) -> None:
        if isinstance(value, str):
            values.append(value)
        elif isinstance(value, list):
            for child in value:
                collect(child)
        elif isinstance(value, dict):
            for child in value.values():
                collect(child)

    collect(json.loads(_read_encoded_text(path)))
    return "\n".join(values)


def _read_docx(path: str) -> str:
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(path)
    lines: list[str] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            value = Paragraph(child, document).text
            if value.strip():
                lines.append(value)
        elif isinstance(child, CT_Tbl):
            for row in Table(child, document).rows:
                values = [cell.text.replace("\n", " ") for cell in row.cells]
                if any(value.strip() for value in values):
                    lines.append("\t".join(values))
    for section in document.sections:
        for region in (section.header, section.footer):
            lines.extend(item.text for item in region.paragraphs if item.text.strip())
    return "\n".join(lines)


def _read_pdf(path: str) -> str:
    import fitz

    document = fitz.open(path)
    try:
        if document.page_count > MAX_PDF_PAGES:
            raise ValueError(f"PDF 页数超过安全上限（{MAX_PDF_PAGES:,} 页）。")
        text = "\n\n".join(page.get_text("text", sort=True) for page in document)
    finally:
        document.close()
    if not text.strip():
        raise ValueError("PDF 没有可读取的文本层；扫描版 PDF 请先进行文字识别。")
    return text


def _read_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = []
    cell_count = 0
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows(values_only=True):
                cell_count += len(row)
                if cell_count > MAX_SPREADSHEET_CELLS:
                    raise ValueError(
                        f"表格单元格数量超过安全上限（{MAX_SPREADSHEET_CELLS:,} 个）。"
                    )
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    lines.append("\t".join(values))
    finally:
        workbook.close()
    return "\n".join(lines)


def _read_pptx(path: str) -> str:
    from pptx import Presentation

    lines: list[str] = []
    presentation = Presentation(path)
    if len(presentation.slides) > MAX_PRESENTATION_SLIDES:
        raise ValueError(
            f"演示文稿页数超过安全上限（{MAX_PRESENTATION_SLIDES:,} 页）。"
        )
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                lines.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        lines.append("\t".join(values))
    return "\n".join(lines)


def _read_open_document(path: str) -> str:
    with zipfile.ZipFile(path) as archive:
        source = archive.read("content.xml")
    root = ET.fromstring(source)
    lines = [value.strip() for value in root.itertext() if value.strip()]
    return "\n".join(lines)


def _read_legacy_office(path: str, extension: str) -> tuple[str, str]:
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        raise ValueError("读取该格式需要安装本机 WPS 或 Microsoft Office。") from exc

    if extension in {".doc", ".wps", ".wpt"}:
        application_names = ("Kwps.Application", "Word.Application")
        document_kind = "文字文档"
        reader = _read_legacy_word
    elif extension in {".xls", ".et", ".ett"}:
        application_names = ("Ket.Application", "Excel.Application")
        document_kind = "表格文档"
        reader = _read_legacy_sheet
    else:
        application_names = ("Kwpp.Application", "PowerPoint.Application")
        document_kind = "演示文稿"
        reader = _read_legacy_presentation

    errors: list[str] = []
    pythoncom.CoInitialize()
    try:
        for application_name in application_names:
            application = None
            try:
                application = win32com.client.DispatchEx(application_name)
                application.Visible = False
                return reader(application, path), f"已通过本机 WPS/Office 只读提取{document_kind}"
            except Exception as exc:
                errors.append(f"{application_name}: {exc}")
            finally:
                if application is not None:
                    try:
                        application.Quit()
                    except Exception:
                        pass
    finally:
        pythoncom.CoUninitialize()
    raise ValueError(f"无法调用本机 WPS 或 Office 读取{document_kind}。")


def _read_legacy_word(application: Any, path: str) -> str:
    document = None
    try:
        try:
            document = application.Documents.Open(
                os.path.abspath(path),
                ConfirmConversions=False,
                ReadOnly=True,
                AddToRecentFiles=False,
                Visible=False,
            )
        except Exception:
            document = application.Documents.Open(os.path.abspath(path), False, True)
        text = str(document.Content.Text)
        if len(text) > MAX_EXTRACTED_CHARACTERS:
            raise ValueError("文字文档内容超过安全上限。")
        return text
    finally:
        if document is not None:
            document.Close(False)


def _read_legacy_sheet(application: Any, path: str) -> str:
    workbook = None
    try:
        workbook = application.Workbooks.Open(os.path.abspath(path), ReadOnly=True)
        lines: list[str] = []
        cell_count = 0
        for worksheet in workbook.Worksheets:
            used_range = worksheet.UsedRange
            try:
                cell_count += int(used_range.Rows.Count) * int(used_range.Columns.Count)
            except (AttributeError, TypeError, ValueError):
                pass
            if cell_count > MAX_SPREADSHEET_CELLS:
                raise ValueError("表格单元格数量超过安全上限。")
            values = used_range.Value
            rows = values if isinstance(values, tuple) else ((values,),)
            for row in rows:
                cells = row if isinstance(row, tuple) else (row,)
                text = ["" if value is None else str(value) for value in cells]
                if any(text):
                    lines.append("\t".join(text))
        return "\n".join(lines)
    finally:
        if workbook is not None:
            workbook.Close(False)


def _read_legacy_presentation(application: Any, path: str) -> str:
    presentation = None
    try:
        presentation = application.Presentations.Open(os.path.abspath(path), WithWindow=False)
        if int(presentation.Slides.Count) > MAX_PRESENTATION_SLIDES:
            raise ValueError("演示文稿页数超过安全上限。")
        lines: list[str] = []
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                try:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        value = str(shape.TextFrame.TextRange.Text).strip()
                        if value:
                            lines.append(value)
                except Exception:
                    continue
        return "\n".join(lines)
    finally:
        if presentation is not None:
            presentation.Close()
